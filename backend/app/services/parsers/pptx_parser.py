from pathlib import Path

from pptx import Presentation


def extract_text_from_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)
